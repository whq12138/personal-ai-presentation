"""
PPTX Export Service — converts Slide-JSON into a downloadable .pptx file
using python-pptx. Each slide layout maps to specific PowerPoint shapes.
"""

import io
import logging
from typing import Optional

from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from app.models.slide import (
    Presentation,
    Slide,
    SlideLayout,
    ContentBlock,
    ContentBlockType,
)

logger = logging.getLogger(__name__)

# Layout constants for a 16:9 slide (10 x 5.625 inches)
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)
MARGIN = Inches(0.5)

# Color palette for PPTX (matching our dark theme, inverted for print)
COLORS = {
    "bg_primary": RGBColor(0x0A, 0x19, 0x2F),
    "text_primary": RGBColor(0xE6, 0xF1, 0xFF),
    "text_secondary": RGBColor(0x88, 0x92, 0xB0),
    "accent": RGBColor(0x64, 0xFF, 0xDA),
    "accent_dim": RGBColor(0x1A, 0x3D, 0x3A),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "black": RGBColor(0x00, 0x00, 0x00),
}


def generate_pptx(presentation: Presentation) -> bytes:
    """
    Convert a Presentation object to a .pptx byte stream.

    Each SlideLayout maps to:
    - title → centered title + subtitle
    - two-column → two side-by-side text boxes
    - highlight-number → large number + label + bullet points
    - table → PowerPoint table
    - bullet-list → title + bullet list
    """
    prs = PptxPresentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Use a blank layout for maximum control
    blank_layout = prs.slide_layouts[6]  # blank layout

    for slide in presentation.slides:
        ppt_slide = prs.slides.add_slide(blank_layout)
        _set_slide_background(ppt_slide, slide)

        layout_handlers = {
            SlideLayout.TITLE: _render_title_slide,
            SlideLayout.TWO_COLUMN: _render_two_column_slide,
            SlideLayout.HIGHLIGHT_NUMBER: _render_highlight_number_slide,
            SlideLayout.TABLE: _render_table_slide,
            SlideLayout.BULLET_LIST: _render_bullet_list_slide,
        }

        handler = layout_handlers.get(slide.layout)
        if handler:
            handler(ppt_slide, slide)
        else:
            _render_title_slide(ppt_slide, slide)  # fallback

    # Save to bytes buffer
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()


def _set_slide_background(slide, slide_data: Slide):
    """Set the slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    # Use a dark navy for a professional look (print-friendly dark bg)
    fill.fore_color.rgb = COLORS["bg_primary"]


def _add_title_box(slide, title: str, top: Emu = Inches(0.3), font_size: Pt = Pt(36)):
    """Add a title text box at the top of a slide."""
    left = MARGIN
    width = SLIDE_W - MARGIN * 2
    height = Inches(1.0)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = font_size
    p.font.bold = True
    p.font.color.rgb = COLORS["accent"]
    p.alignment = PP_ALIGN.LEFT
    return tf


def _render_title_slide(slide, slide_data: Slide):
    """Render a title slide: centered title + subtitle + body."""
    title = slide_data.title or "Untitled"
    subtitle = slide_data.subtitle or ""

    # Title
    txBox = slide.shapes.add_textbox(
        Inches(1), Inches(1.5), Inches(8), Inches(1.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS["accent"]
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        txBox2 = slide.shapes.add_textbox(
            Inches(1.5), Inches(3.2), Inches(7), Inches(1.0)
        )
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = COLORS["text_secondary"]
        p2.alignment = PP_ALIGN.CENTER

    # Body
    body = slide_data.body
    if body:
        txBox3 = slide.shapes.add_textbox(
            Inches(2), Inches(4.0), Inches(6), Inches(1.0)
        )
        tf3 = txBox3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = body if isinstance(body, str) else body[0]
        p3.font.size = Pt(14)
        p3.font.color.rgb = COLORS["text_secondary"]
        p3.alignment = PP_ALIGN.CENTER


def _render_two_column_slide(slide, slide_data: Slide):
    """Render a two-column slide."""
    _add_title_box(slide, slide_data.title or "Section")

    col_width = Inches(4.25)
    col_height = Inches(3.5)
    top = Inches(1.5)

    # Left column
    if slide_data.columns and slide_data.columns.left:
        _render_content_blocks(
            slide, slide_data.columns.left,
            left=MARGIN, top=top, width=col_width, height=col_height
        )

    # Right column
    if slide_data.columns and slide_data.columns.right:
        _render_content_blocks(
            slide, slide_data.columns.right,
            left=Inches(5.25), top=top, width=col_width, height=col_height
        )


def _render_highlight_number_slide(slide, slide_data: Slide):
    """Render a highlight-number slide: big number + label + bullets."""
    _add_title_box(slide, slide_data.title or "Key Metric")

    hn = slide_data.highlightNumber
    if hn:
        # Big number
        txBox = slide.shapes.add_textbox(
            Inches(1), Inches(1.5), Inches(3.5), Inches(2.0)
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = hn.value
        p.font.size = Pt(72)
        p.font.bold = True
        p.font.color.rgb = COLORS["accent"]
        p.alignment = PP_ALIGN.CENTER

        # Label
        txBox2 = slide.shapes.add_textbox(
            Inches(1), Inches(3.5), Inches(3.5), Inches(0.8)
        )
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = hn.label
        p2.font.size = Pt(16)
        p2.font.color.rgb = COLORS["text_secondary"]
        p2.alignment = PP_ALIGN.CENTER

        # Suffix
        if hn.suffix:
            txBox3 = slide.shapes.add_textbox(
                Inches(1), Inches(4.0), Inches(3.5), Inches(0.5)
            )
            tf3 = txBox3.text_frame
            p3 = tf3.paragraphs[0]
            p3.text = hn.suffix
            p3.font.size = Pt(14)
            p3.font.color.rgb = COLORS["accent"]
            p3.alignment = PP_ALIGN.CENTER

    # Supporting bullets
    body = slide_data.body
    if body:
        items = body if isinstance(body, list) else [body]
        txBox4 = slide.shapes.add_textbox(
            Inches(5), Inches(1.5), Inches(4.5), Inches(3.5)
        )
        tf4 = txBox4.text_frame
        for i, item in enumerate(items):
            p = tf4.paragraphs[0] if i == 0 else tf4.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = COLORS["text_primary"]
            p.space_after = Pt(8)


def _render_table_slide(slide, slide_data: Slide):
    """Render a table slide."""
    _add_title_box(slide, slide_data.title or "Data Table")

    tbl_data = slide_data.table
    if not tbl_data:
        return

    headers = tbl_data.headers
    rows = tbl_data.rows
    num_rows = len(rows) + 1  # +1 for header
    num_cols = len(headers)

    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(1), Inches(1.8),
        Inches(8), Inches(0.5 * num_rows)
    )
    table = table_shape.table

    # Headers
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.bold = True
            paragraph.font.color.rgb = COLORS["accent"]
            paragraph.alignment = PP_ALIGN.CENTER

    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = COLORS["text_primary"]
                paragraph.alignment = PP_ALIGN.CENTER


def _render_bullet_list_slide(slide, slide_data: Slide):
    """Render a bullet-list slide."""
    _add_title_box(slide, slide_data.title or "Key Points")

    body = slide_data.body
    items = body if isinstance(body, list) else ([body] if body else [])

    txBox = slide.shapes.add_textbox(
        Inches(1), Inches(1.5), Inches(8), Inches(3.5)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS["text_primary"]
        p.space_after = Pt(12)


def _render_content_blocks(slide, blocks: list[ContentBlock], left, top, width, height):
    """Render a list of ContentBlocks (heading/paragraph/image/list) into a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    for block in blocks:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False

        if block.type == ContentBlockType.HEADING:
            p.text = block.text or ""
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = COLORS["accent"]
        elif block.type == ContentBlockType.LIST and block.items:
            for j, item in enumerate(block.items):
                lp = p if j == 0 else tf.add_paragraph()
                lp.text = f"• {item}"
                lp.font.size = Pt(14)
                lp.font.color.rgb = COLORS["text_primary"]
                lp.space_after = Pt(6)
        else:
            p.text = block.text or ""
            p.font.size = Pt(14)
            p.font.color.rgb = COLORS["text_primary"]

        p.space_after = Pt(8)
